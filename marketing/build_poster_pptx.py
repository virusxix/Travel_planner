"""HiddenStay A1 poster — HAWKER-style layout (JCU capstone)."""
from pptx import Presentation
from pptx.util import Mm, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

OUT = "HiddenStay-Poster-HAWKER.pptx"

# HAWKER-style warm browns
CREAM = RGBColor(0xF5, 0xED, 0xE0)
PANEL = RGBColor(0xE8, 0xDC, 0xC8)
TAN = RGBColor(0xD4, 0xC4, 0xA8)
BROWN = RGBColor(0x5C, 0x40, 0x33)
DARK = RGBColor(0x3D, 0x2E, 0x24)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
STAT_RED = RGBColor(0x8B, 0x3A, 0x2A)
GREEN = RGBColor(0x2D, 0x6A, 0x5A)
MUTED = RGBColor(0x6B, 0x5A, 0x48)
SWOT_S = RGBColor(0x5B, 0x7C, 0x99)
SWOT_W = RGBColor(0xD4, 0xA8, 0x43)
SWOT_O = RGBColor(0xD4, 0x84, 0x5B)
SWOT_T = RGBColor(0x8B, 0x6B, 0x9E)
ORANGE = RGBColor(0xC4, 0x7A, 0x3A)
YELLOW = RGBColor(0xE8, 0xC5, 0x47)


def rect(slide, l, t, w, h, fill, line=BROWN, lw=1.5, radius=None):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(kind, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(lw)
    else:
        s.line.fill.background()
    return s


def txt(slide, l, t, w, h, text, size=14, bold=False, color=DARK, align=PP_ALIGN.LEFT, font="Arial"):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    r = p.runs[0]
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = font
    return tb


def section_hdr(slide, l, t, w, label):
    txt(slide, l, t, w, Mm(10), label, size=22, bold=True, color=BROWN, font="Arial Black")


def pill(slide, l, t, w, h, label, fill=BROWN, fg=WHITE, size=11):
    rect(slide, l, t, w, h, fill, line=fill)
    txt(slide, l, t, w, h, label, size=size, bold=True, color=fg, align=PP_ALIGN.CENTER)


def flow_node(slide, l, t, w, h, text, fill=TAN, size=10):
    rect(slide, l, t, w, h, fill, radius=True)
    txt(slide, l + Mm(1), t + Mm(1.5), w - Mm(2), h - Mm(2), text, size=size, bold=True, color=DARK, align=PP_ALIGN.CENTER)


def main():
    prs = Presentation()
    prs.slide_width = Mm(594)
    prs.slide_height = Mm(841)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    rect(slide, 0, 0, Mm(594), Mm(841), CREAM, line=None, lw=0)

    m = Mm(38)
    cw = Mm(518)
    gap = Mm(6)
    half = (cw - gap) / 2
    xL, xR = m, m + half + gap

    # ── HEADER (HAWKER: logos + big title + subtitle) ──
    y = m
    rect(slide, xL, y, Mm(48), Mm(24), CREAM, line=MUTED, lw=1)
    txt(slide, xL, y + Mm(5), Mm(48), Mm(14), "JCUS 18\nSEAL", size=10, bold=True, color=MUTED, align=PP_ALIGN.CENTER)

    cx = xL + Mm(52)
    cw_title = cw - Mm(114)
    txt(slide, cx, y, cw_title, Mm(16), "HiddenStay", size=60, bold=True, color=BROWN, align=PP_ALIGN.CENTER, font="Times New Roman")
    txt(slide, cx, y + Mm(17), cw_title, Mm(8), "FAIR COMMISSION HOMESTAYS + AI TRIP PLANNER · SOUTHEAST ASIA", size=13, bold=True, color=MUTED, align=PP_ALIGN.CENTER)

    rect(slide, xL + cw - Mm(62), y, Mm(62), Mm(24), CREAM, line=MUTED, lw=1)
    txt(slide, xL + cw - Mm(62), y + Mm(4), Mm(62), Mm(16), "JAMES COOK\nUNIVERSITY\nSINGAPORE", size=9, bold=True, color=MUTED, align=PP_ALIGN.CENTER)

    y += Mm(30)

    # ═══════════════════════════════════════
    # LEFT TOP — THE ISSUE (big stats like HAWKER "121" "13,430")
    # ═══════════════════════════════════════
    issue_h = Mm(95)
    rect(slide, xL, y, half, issue_h, PANEL)
    section_hdr(slide, xL + Mm(5), y + Mm(4), half - Mm(10), "THE ISSUE")

    txt(slide, xL + Mm(5), y + Mm(16), half - Mm(10), Mm(14),
        "Independent homestays lose 15–30% per booking to OTAs. Travellers plan on Google and book on different apps — information and booking are split.",
        size=13, color=DARK)

    # Big stat row (HAWKER style)
    sy = y + Mm(34)
    sw = (half - Mm(16)) / 3
    stats = [("15–30%", "OTA FEE TODAY"), ("5%", "HIDDENSTAY FEE"), ("95%", "HOST KEEPS")]
    for i, (num, lbl) in enumerate(stats):
        bx = xL + Mm(5) + i * (sw + Mm(3))
        txt(slide, bx, sy, sw, Mm(14), num, size=40, bold=True, color=STAT_RED, align=PP_ALIGN.CENTER, font="Arial Black")
        txt(slide, bx, sy + Mm(14), sw, Mm(8), lbl, size=11, bold=True, color=MUTED, align=PP_ALIGN.CENTER)

    txt(slide, xL + Mm(5), y + issue_h - Mm(12), half - Mm(10), Mm(10),
        "SGD 100 stay → host saves SGD 13+ vs 18% OTA.", size=12, bold=True, color=DARK)

    # ═══════════════════════════════════════
    # RIGHT TOP — BUSINESS MODEL (flowchart like HAWKER)
    # ═══════════════════════════════════════
    rect(slide, xR, y, half, issue_h, PANEL)
    section_hdr(slide, xR + Mm(5), y + Mm(4), half - Mm(10), "BUSINESS MODEL")

    nw, nh = Mm(38), Mm(14)
    top_y = y + Mm(18)
    nodes = [
        (xR + Mm(8), "FOCUS", "Supply +\nplanner"),
        (xR + half / 2 - nw / 2, "SUBSCRIPTION", "Free\nStarter · 5%"),
        (xR + half - nw - Mm(8), "REVENUE", "4 income\nstreams"),
    ]
    for nx, hdr, body in nodes:
        flow_node(slide, nx, top_y, nw, nh, hdr + "\n" + body, fill=TAN, size=9)

    # Revenue chain (arrows like HAWKER)
    chain_y = top_y + nh + Mm(6)
    chain_items = ["5% booking fee", "Growth SGD 50/mo", "Featured SGD 99", "Ads 7–10% (15% cap)"]
    cw2 = half - Mm(16)
    for i, item in enumerate(chain_items):
        cy = chain_y + i * Mm(11)
        if i > 0:
            txt(slide, xR + Mm(5), cy - Mm(5), Mm(8), Mm(5), "↓", size=14, bold=True, color=BROWN, align=PP_ALIGN.CENTER)
        flow_node(slide, xR + Mm(12), cy, cw2 - Mm(8), Mm(9), item, fill=ORANGE if i == 0 else YELLOW if i == 1 else TAN, size=10)

    y += issue_h + gap

    # ═══════════════════════════════════════
    # LEFT MID — WHAT IF + MISSION/VISION (HAWKER hook)
    # ═══════════════════════════════════════
    mid_h = Mm(72)
    rect(slide, xL, y, half, mid_h, PANEL)

    hook_w = half - Mm(50)
    txt(slide, xL + Mm(5), y + Mm(6), hook_w, Mm(20),
        "WHAT IF HOSTS COULD GET IT ALL IN ONE APP?", size=16, bold=True, color=DARK, font="Arial Black")
    txt(slide, xL + Mm(5), y + Mm(22), hook_w, Mm(10),
        "List free · Keep 95% · Plan & book in one place.", size=13, color=MUTED)

    # App logo circle (HAWKER hawker logo)
    logo = slide.shapes.add_shape(MSO_SHAPE.OVAL, xL + half - Mm(22), y + Mm(8), Mm(18), Mm(18))
    logo.fill.solid()
    logo.fill.fore_color.rgb = GREEN
    logo.line.color.rgb = BROWN
    logo.line.width = Pt(2)
    txt(slide, xL + half - Mm(22), y + Mm(12), Mm(18), Mm(10), "HS", size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Mission / Vision pills
    mv_y = y + Mm(38)
    mv_w = (half - Mm(14)) / 2
    for i, (title, body) in enumerate([
        ("MISSION", "Fair host economics + seamless plan → map → book for travellers."),
        ("VISION", "Trusted direct-booking for hidden stays across SEA by 2030."),
    ]):
        mx = xL + Mm(5) + i * (mv_w + Mm(4))
        rect(slide, mx, mv_y, mv_w, Mm(28), TAN, lw=1)
        pill(slide, mx + Mm(2), mv_y + Mm(2), mv_w - Mm(4), Mm(7), title, size=10)
        txt(slide, mx + Mm(3), mv_y + Mm(10), mv_w - Mm(6), Mm(16), body, size=11, color=DARK)

    # ═══════════════════════════════════════
    # RIGHT MID — WHY US? + SWOT table + acronym (HAWKER)
    # ═══════════════════════════════════════
    rect(slide, xR, y, half, mid_h, PANEL)
    section_hdr(slide, xR + Mm(5), y + Mm(4), half - Mm(10), "WHY US?")

    # SWOT table (coloured rows like HAWKER)
    sw_x = xR + Mm(5)
    sw_w = Mm(72)
    swot = [
        (SWOT_S, "STRENGTHS", "5% vs 15–30% OTAs\nFree host portal\nPlanner → book loop"),
        (SWOT_W, "WEAKNESSES", "Thin unit economics\nNew brand"),
        (SWOT_O, "OPPORTUNITIES", "SEA homestay supply\nSEO + planner content"),
        (SWOT_T, "THREATS", "OTA dominance\nGuest discovery cost"),
    ]
    row_h = Mm(13)
    for i, (col, title, body) in enumerate(swot):
        ry = y + Mm(16) + i * row_h
        rect(slide, sw_x, ry, sw_w, row_h - Mm(1), col, line=col)
        txt(slide, sw_x + Mm(2), ry + Mm(1), sw_w - Mm(4), Mm(5), title, size=9, bold=True, color=WHITE)
        txt(slide, sw_x + sw_w + Mm(3), ry, half - sw_w - Mm(18), row_h, body, size=10, color=DARK)

    # HIDDENSTAY acronym strip (like H-A-W-K-E-R vertical circles)
    ac_x = xR + half - Mm(28)
    ac_top = y + Mm(16)
    letters = [
        ("H", "Host keeps 95%"),
        ("I", "Independent stays"),
        ("D", "Direct booking"),
        ("D", "Discovery map"),
        ("E", "Earnings portal"),
        ("N", "Night-stay trips"),
        ("S", "Smart planner"),
        ("T", "Transparent 5%"),
    ]
    for i, (letter, desc) in enumerate(letters):
        ly = ac_top + i * Mm(7.5)
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, ac_x, ly, Mm(7), Mm(7))
        circ.fill.solid()
        circ.fill.fore_color.rgb = BROWN
        circ.line.fill.background()
        txt(slide, ac_x, ly + Mm(0.5), Mm(7), Mm(6), letter, size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        if i < 4:  # only show desc for first few to avoid clutter
            txt(slide, ac_x + Mm(9), ly, Mm(18), Mm(7), desc, size=7, color=MUTED)

    y += mid_h + gap

    # ═══════════════════════════════════════
    # BOTTOM — Phones (left) + Features 1–7 (right) — HAWKER layout
    # ═══════════════════════════════════════
    footer_y = Mm(665)
    bot_h = footer_y - y - Mm(3)

    # Phone area
    phone_area_w = half + Mm(20)
    rect(slide, xL, y, phone_area_w, bot_h, TAN, lw=1.5)

    pw, ph = Mm(38), Mm(68)
    total_phones = pw * 3 + Mm(12)
    px0 = xL + (phone_area_w - total_phones) / 2
    py0 = y + (bot_h - ph) / 2
    labels = ["AI PLANNER", "MAP + BOOK", "HOST PORTAL"]
    hints = [
        "Day 1 · Chiang Mai\n● Old town\n● Homestay\nDay 2 · Countryside",
        "[ MAP VIEW ]\n● Stay A  SGD 50/n\n● Stay B  SGD 45/n\n→ Book now",
        "EARNINGS\nAvailable  SGD 285\nPending    SGD 95\n5% fee · 95% yours",
    ]
    for i, (lbl, hint) in enumerate(zip(labels, hints)):
        px = px0 + i * (pw + Mm(6))
        phone = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, px, py0, pw, ph)
        phone.fill.solid()
        phone.fill.fore_color.rgb = DARK
        phone.line.color.rgb = BROWN
        phone.line.width = Pt(2)
        txt(slide, px, py0 + Mm(2), pw, Mm(6), lbl, size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        rect(slide, px + Mm(2), py0 + Mm(9), pw - Mm(4), ph - Mm(11), RGBColor(0x0F, 0x2A, 0x3A), line=RGBColor(0x1A, 0x3D, 0x4A), lw=0.5)
        txt(slide, px + Mm(3), py0 + Mm(11), pw - Mm(6), ph - Mm(14), hint, size=8, color=RGBColor(0xA8, 0xE6, 0xCF), align=PP_ALIGN.CENTER)

    txt(slide, xL + Mm(5), y + Mm(3), phone_area_w - Mm(10), Mm(6),
        "[ Replace phone screens with real app screenshots ]", size=9, bold=True, color=MUTED, align=PP_ALIGN.CENTER)

    # Features column (HAWKER numbered list)
    feat_x = xR
    feat_w = half - Mm(20)
    rect(slide, feat_x, y, feat_w, bot_h, PANEL)
    section_hdr(slide, feat_x + Mm(4), y + Mm(4), feat_w - Mm(8), "FEATURES")

    feats = [
        "5% MARKETPLACE",
        "AI PLANNER + MAP",
        "PLAN → BOOK FLOW",
        "FREE HOST PORTAL",
        "GROWTH TOOLS (SEO)",
        "MULTI-NIGHT TRIPS",
        "18-MO SURVIVAL PLAN",
    ]
    fh = (bot_h - Mm(22)) / 7
    for i, f in enumerate(feats):
        fy = y + Mm(16) + i * (fh + Mm(1.5))
        rect(slide, feat_x + Mm(4), fy, feat_w - Mm(8), fh - Mm(1), TAN, lw=1)
        num = slide.shapes.add_shape(MSO_SHAPE.OVAL, feat_x + Mm(6), fy + fh / 2 - Mm(4), Mm(8), Mm(8))
        num.fill.solid()
        num.fill.fore_color.rgb = BROWN
        num.line.fill.background()
        txt(slide, feat_x + Mm(6), fy + fh / 2 - Mm(3.5), Mm(8), Mm(7), str(i + 1), size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(slide, feat_x + Mm(16), fy + fh / 2 - Mm(4), feat_w - Mm(22), Mm(8), f, size=12, bold=True, color=DARK)

    # ── FOOTER ──
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, xL, footer_y, cw, Mm(1))
    line.fill.solid()
    line.fill.fore_color.rgb = BROWN
    line.line.fill.background()

    txt(slide, xL, footer_y + Mm(4), cw / 2, Mm(16),
        "Group 63 I — Zaw Latt Naung · Oakkar Phyoe · Eka Dian Tara Tiu · Chenyang Gu · Yihua Deng\nBU3102 / CP3102 Capstone · James Cook University Singapore",
        size=11, color=MUTED)
    txt(slide, xL + cw / 2, footer_y + Mm(4), cw / 2, Mm(16),
        "Break-even Month 18 · 55 hosts · 92 bookings/mo\nhiddenstay.app (pilot)",
        size=11, bold=True, color=BROWN, align=PP_ALIGN.RIGHT)

    prs.save(OUT)
    print(f"Saved {OUT}")


if __name__ == "__main__":
    main()
